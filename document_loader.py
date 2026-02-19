import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

# ✅ Unstructured fallback
from unstructured.partition.auto import partition


def load_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_pdf(file_path: str) -> str:
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text("text") + "\n"
    return text.strip()


def load_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs]).strip()


def load_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs).strip()


def load_image_ocr(file_path: str) -> str:
    img = Image.open(file_path)
    return pytesseract.image_to_string(img)


# ✅ Unstructured fallback loader
def load_with_unstructured(file_path: str) -> str:
    elements = partition(filename=file_path)
    text = "\n".join([el.text for el in elements if el.text])
    return text.strip()


def load_document(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    try:
        # ✅ Your optimized loaders first
        if ext == ".txt":
            return load_txt(file_path)

        elif ext == ".pdf":
            text = load_pdf(file_path)
            if text.strip():
                return text
            # If PDF extraction gives blank (scanned), fallback
            return load_with_unstructured(file_path)

        elif ext == ".docx":
            return load_docx(file_path)

        elif ext == ".pptx":
            return load_pptx(file_path)

        elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
            return load_image_ocr(file_path)

        # ✅ Unknown format -> fallback directly
        return load_with_unstructured(file_path)

    except Exception as e:
        # ✅ If anything fails, fallback to unstructured
        try:
            return load_with_unstructured(file_path)
        except Exception:
            raise ValueError(f"Failed to load document: {file_path}. Error: {str(e)}")

